"""Format-aware extraction into chunks with exact, structured source locations."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
import csv
import re
import tempfile
from time import monotonic

from app.config import settings
from app.services.chunking import chunk_text
from app.services.document_loader import DocumentParseError, extract_text
from app.services.image_processor import (
    IMAGE_EXTENSIONS,
    ImageProcessingError,
    chunk_image_text,
    extract_image_text,
)
from app.services.pdf_layout import extract_pdf_page_texts

DOCX_PARAGRAPH_CHUNK_WORDS = 800


@dataclass(frozen=True)
class SourceChunk:
    text: str
    source_type: str
    location: dict[str, object]


def validate_source_location(source_type: str, location: dict[str, object]) -> None:
    """Reject chunks that cannot produce a precise citation for their source type."""
    required: dict[str, tuple[str, ...]] = {
        "pdf": ("page_start", "page_end"),
        "powerpoint": ("slide_start", "slide_end", "content_type"),
        "excel": (
            "sheet_name", "row_start", "row_end", "column_start",
            "column_end", "cell_range", "hidden_sheet",
        ),
        "csv": ("row_start", "row_end"),
        "text": ("line_start", "line_end"),
    }
    missing = [
        key for key in required.get(source_type, ())
        if key not in location or location[key] is None
    ]
    if missing:
        raise DocumentParseError(
            f"Invalid {source_type} source location; missing {', '.join(missing)}."
        )
    for start, end in (
        ("page_start", "page_end"),
        ("slide_start", "slide_end"),
        ("row_start", "row_end"),
        ("line_start", "line_end"),
    ):
        if start in location and end in location:
            if int(location[start]) < 1 or int(location[end]) < int(location[start]):
                raise DocumentParseError(
                    f"Invalid {source_type} source range: {start}/{end}."
                )


def _split(text: str, source_type: str, location: dict[str, object]) -> list[SourceChunk]:
    return [
        SourceChunk(value, source_type, {**location, "part": index})
        for index, value in enumerate(chunk_text(text), start=1)
    ]


def _table_row_text(headers: list[str], cells: list[str], prefix: str, row_number: int) -> str:
    """Keep table cells bound to their headers instead of flattening rows into prose."""
    labels = [header.strip() or f"Column {index + 1}" for index, header in enumerate(headers)]
    if row_number == 1:
        return f"{prefix} | Headers: " + " | ".join(labels)
    return f"{prefix} | " + " | ".join(
        f"{labels[index] if index < len(labels) else f'Column {index + 1}'}: {value}"
        for index, value in enumerate(cells)
        if value
    )


def _ocr_text_is_duplicate(ocr_text: str, native_text: str) -> bool:
    """Avoid indexing OCR text that repeats already extracted native text."""
    ocr_tokens = set(re.findall(r"[a-z0-9]+", ocr_text.casefold()))
    native_tokens = set(re.findall(r"[a-z0-9]+", native_text.casefold()))
    if len(ocr_tokens) < 4 or not native_tokens:
        return False
    return len(ocr_tokens & native_tokens) / len(ocr_tokens) >= 0.80


def _safe_image_suffix(name: str | None, content_type: str | None = None) -> str:
    """Choose a Pillow-friendly extension without trusting archive names."""
    suffix = Path(name or "").suffix.lower()
    if suffix in IMAGE_EXTENSIONS:
        return suffix
    if content_type:
        guessed = "." + content_type.split("/")[-1].lower().replace("jpeg", "jpg")
        if guessed in IMAGE_EXTENSIONS:
            return guessed
    return ".png"


def _image_chunks_from_bytes(
    *,
    blob: bytes,
    filename_hint: str,
    source_type: str,
    location: dict[str, object],
    native_text: str = "",
) -> list[SourceChunk]:
    """OCR one embedded image through the same bounded image pipeline used for uploads."""
    if not blob or len(blob) > settings.embedded_ocr_max_decoded_bytes:
        return []
    try:
        from PIL import Image

        with Image.open(BytesIO(blob)) as image:
            width, height = image.size
            frames = int(getattr(image, "n_frames", 1))
            if width * height * max(frames, 1) > settings.embedded_ocr_max_pixels:
                return []
    except Exception:
        return []
    suffix = _safe_image_suffix(filename_hint)
    try:
        with tempfile.TemporaryDirectory(prefix="rag-embedded-ocr-") as temporary:
            path = Path(temporary) / f"embedded{suffix}"
            path.write_bytes(blob)
            text = extract_image_text(path)
    except ImageProcessingError as error:
        if error.code in {"ocr_unavailable", "ocr_timeout", "ocr_processing_failed"}:
            raise DocumentParseError(str(error), code=error.code) from error
        return []
    except Exception:
        # Embedded OCR is best-effort for malformed image data, not infrastructure failures.
        return []
    if _ocr_text_is_duplicate(text, native_text):
        return []
    return [
        SourceChunk(value, source_type, {**location, "part": index})
        for index, value in enumerate(chunk_image_text(text), start=1)
    ]


def _docx_paragraph_chunks(document) -> list[SourceChunk]:
    """Batch short DOCX paragraphs so large prose documents do not trip chunk-count limits."""
    result: list[SourceChunk] = []
    buffer: list[str] = []
    buffer_words = 0
    paragraph_start: int | None = None
    paragraph_end: int | None = None
    section_index = 1
    section_start = section_index
    styles: set[str] = set()

    def flush() -> None:
        nonlocal buffer, buffer_words, paragraph_start, paragraph_end, section_start, styles
        if not buffer or paragraph_start is None or paragraph_end is None:
            return
        result.append(SourceChunk(
            "\n\n".join(buffer),
            "word",
            {
                "section_number": section_start,
                "paragraph_start": paragraph_start,
                "paragraph_end": paragraph_end,
                "styles": sorted(styles),
                "content_type": "prose",
            },
        ))
        buffer = []
        buffer_words = 0
        paragraph_start = None
        paragraph_end = None
        section_start = section_index
        styles = set()

    for paragraph_index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if paragraph._p.xpath("./w:pPr/w:sectPr"):
            flush()
            section_index += 1
        if not text:
            continue
        words = len(text.split())
        style = paragraph.style.name if paragraph.style else None
        if style and style.casefold().startswith("heading"):
            flush()
        if words > DOCX_PARAGRAPH_CHUNK_WORDS:
            flush()
            result.extend(_split(text, "word", {
                "section_number": section_index,
                "paragraph_start": paragraph_index,
                "paragraph_end": paragraph_index,
                "style": style,
                "content_type": "prose",
            }))
            continue
        if buffer and buffer_words + words > DOCX_PARAGRAPH_CHUNK_WORDS:
            flush()
        if paragraph_start is None:
            paragraph_start = paragraph_index
            section_start = section_index
        paragraph_end = paragraph_index
        buffer.append(text)
        buffer_words += words
        if style:
            styles.add(style)
    flush()
    return result


def _docx_image_locations(document) -> dict[str, dict[str, object]]:
    """Map embedded image relationship IDs to nearby paragraph provenance."""
    locations: dict[str, dict[str, object]] = {}
    section_index = 1
    current_heading = ""
    for paragraph_index, paragraph in enumerate(document.paragraphs, start=1):
        if paragraph._p.xpath("./w:pPr/w:sectPr"):
            section_index += 1
        style = paragraph.style.name if paragraph.style else ""
        if style.casefold().startswith("heading") and paragraph.text.strip():
            current_heading = paragraph.text.strip()
        for rel_id in re.findall(r'r:embed="([^"]+)"', paragraph._p.xml):
            locations.setdefault(rel_id, {
                "section_number": section_index,
                "paragraph_start": paragraph_index,
                "paragraph_end": paragraph_index,
                "heading": current_heading,
                "nearby_text": paragraph.text.strip(),
            })
    return locations


def _docx_image_chunks(document) -> list[SourceChunk]:
    """OCR internal DOCX media parts without following links or active content."""
    if settings.embedded_ocr_max_images_per_document == 0:
        return []
    result: list[SourceChunk] = []
    locations = _docx_image_locations(document)
    image_index = 0
    for rel_id, part in sorted(document.part.related_parts.items()):
        content_type = str(getattr(part, "content_type", ""))
        if not content_type.startswith("image/"):
            continue
        image_index += 1
        if image_index > settings.embedded_ocr_max_images_per_document:
            break
        nearby = locations.get(rel_id, {})
        location = {
            **{key: value for key, value in nearby.items() if key != "nearby_text"},
            "relationship_id": rel_id,
            "image_index": image_index,
            "content_type": "image_ocr",
        }
        result.extend(_image_chunks_from_bytes(
            blob=bytes(getattr(part, "blob", b"")),
            filename_hint=str(getattr(part, "partname", "")),
            source_type="word",
            location=location,
            native_text=str(nearby.get("nearby_text") or ""),
        ))
    return result


def _pdf_table_cells(line: str) -> list[str]:
    """Recognize layout-preserved PDF table rows without treating prose as columns."""
    value = line.strip().strip("|")
    if "|" in value:
        cells = [cell.strip() for cell in value.split("|")]
    elif "\t" in value:
        cells = [cell.strip() for cell in value.split("\t")]
    else:
        cells = [cell.strip() for cell in re.split(r"\s{2,}", value)]
    cells = [cell for cell in cells if cell]
    return cells if len(cells) >= 2 else []


def _pdf_table_chunks(page_number: int, text: str) -> tuple[list[SourceChunk], set[int]]:
    """Extract complete PDF table rows and report lines removed from prose chunking."""
    lines = text.splitlines()
    result: list[SourceChunk] = []
    consumed: set[int] = set()
    table_number = 0
    cursor = 0
    while cursor < len(lines):
        cells = _pdf_table_cells(lines[cursor])
        if not cells:
            cursor += 1
            continue
        start = cursor
        rows: list[tuple[int, list[str]]] = []
        while cursor < len(lines):
            row_cells = _pdf_table_cells(lines[cursor])
            if not row_cells or len(row_cells) != len(cells):
                break
            rows.append((cursor, row_cells))
            cursor += 1
        if len(rows) < 2:
            cursor = start + 1
            continue
        table_number += 1
        headers = rows[0][1]
        for table_row_index, (position, row_cells) in enumerate(rows, start=1):
            row_number = position + 1
            result.append(SourceChunk(
                _table_row_text(headers, row_cells, f"PDF table {table_number}", table_row_index),
                "pdf",
                {
                    "page_start": page_number,
                    "page_end": page_number,
                    "table_number": table_number,
                    "row_start": row_number,
                    "row_end": row_number,
                    "header_context": headers,
                    "content_type": "table",
                },
            ))
            consumed.add(position)
    return result, consumed


def _pdf(path: Path) -> list[SourceChunk]:
    try:
        result: list[SourceChunk] = []
        native_by_page: dict[int, str] = {}
        for page_text in extract_pdf_page_texts(path):
            native_by_page[page_text.page_number] = page_text.text
            table_chunks, table_lines = _pdf_table_chunks(page_text.page_number, page_text.text)
            result.extend(table_chunks)
            prose = "\n".join(
                line for index, line in enumerate(page_text.text.splitlines())
                if index not in table_lines
            )
            for block_index, value in enumerate(chunk_text(prose), start=1):
                result.append(SourceChunk(value, "pdf", {
                    "page_start": page_text.page_number,
                    "page_end": page_text.page_number,
                    "block_ids": [f"p{page_text.page_number}-b{block_index}"],
                    "bounding_boxes": [],
                    "part": block_index,
                }))
        result.extend(_pdf_image_chunks(path, native_by_page))
        return result
    except DocumentParseError:
        raise
    except Exception as error:
        raise DocumentParseError("The PDF file could not be read.") from error


def _pdf_image_chunks(path: Path, native_by_page: dict[int, str]) -> list[SourceChunk]:
    """OCR bounded internal PDF image XObjects with page provenance."""
    if settings.embedded_ocr_max_images_per_document == 0:
        return []
    try:
        from pypdf import PdfReader
    except Exception:
        return []
    result: list[SourceChunk] = []
    image_count = 0
    try:
        pages = PdfReader(str(path)).pages
        for page_number, page in enumerate(pages, start=1):
            page_images = list(getattr(page, "images", []) or [])
            for page_image_index, image in enumerate(page_images[:settings.embedded_ocr_max_images_per_page], start=1):
                image_count += 1
                if image_count > settings.embedded_ocr_max_images_per_document:
                    return result
                blob = bytes(getattr(image, "data", b"") or b"")
                name = str(getattr(image, "name", "") or f"page-{page_number}-image-{page_image_index}.png")
                result.extend(_image_chunks_from_bytes(
                    blob=blob,
                    filename_hint=name,
                    source_type="pdf",
                    location={
                        "page_start": page_number,
                        "page_end": page_number,
                        "image_index": image_count,
                        "page_image_index": page_image_index,
                        "content_type": "image_ocr",
                    },
                    native_text=native_by_page.get(page_number, ""),
                ))
    except DocumentParseError:
        raise
    except Exception:
        return result
    return result


def _is_nonempty_cell(value: object) -> bool:
    """Treat blank strings like empty cells when detecting workbook headers."""
    return value is not None and (not isinstance(value, str) or bool(value.strip()))


def _vertical_merged_values(worksheet) -> dict[tuple[int, int], object]:
    """Preserve labels from vertically merged header cells without spreading titles."""
    values: dict[tuple[int, int], object] = {}
    for cell_range in worksheet.merged_cells.ranges:
        if cell_range.min_col != cell_range.max_col:
            continue
        value = worksheet.cell(cell_range.min_row, cell_range.min_col).value
        for row_number in range(cell_range.min_row + 1, cell_range.max_row + 1):
            values[(row_number, cell_range.min_col)] = value
    return values


def _header_row_number(rows: list[tuple[int, list[object]]]) -> int | None:
    """Choose the most table-like leading row for field labels."""
    best: tuple[int, int] | None = None
    for index, (row_number, values) in enumerate(rows[:20]):
        present = [value for value in values if _is_nonempty_cell(value)]
        if not present:
            continue
        text_count = sum(isinstance(value, str) for value in present)
        distinct = len({str(value).strip().casefold() for value in present})
        following_width = (
            sum(_is_nonempty_cell(value) for value in rows[index + 1][1])
            if index + 1 < len(rows) else 0
        )
        score = len(present) * 3 + text_count * 2 + distinct + min(following_width, len(present))
        if len(present) == 1 and len(rows) > index + 1:
            score -= 4
        if best is None or score > best[0]:
            best = (score, row_number)
    return best[1] if best else None


def _deduplicated_headers(cells: list[tuple[int, object]]) -> dict[int, str]:
    """Give repeated worksheet headers stable distinct labels for source chunks."""
    headers: dict[int, str] = {}
    used: dict[str, int] = {}
    for column, raw_header in cells:
        if not _is_nonempty_cell(raw_header):
            continue
        header = str(raw_header).strip()
        key = header.casefold()
        used[key] = used.get(key, 0) + 1
        headers[column] = header if used[key] == 1 else f"{header} ({used[key]})"
    return headers


def _pptx(path: Path) -> list[SourceChunk]:
    try:
        from pptx import Presentation

        result: list[SourceChunk] = []
        slides = Presentation(str(path)).slides
        if len(slides) > settings.max_powerpoint_slides:
            raise DocumentParseError("Presentation exceeds the configured slide limit.")
        for slide_number, slide in enumerate(slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                shape_name = str(getattr(shape, "name", f"Shape {shape_index}"))
                shape_id = f"shape-{shape_index}"
                if shape == getattr(slide.shapes, "title", None):
                    shape_type = "title"
                elif getattr(shape, "has_table", False):
                    shape_type = "table"
                elif getattr(shape, "has_chart", False):
                    shape_type = "chart"
                else:
                    shape_type = "text_box"
                common = {
                    "slide_start": slide_number,
                    "slide_end": slide_number,
                    "slide_number": slide_number,
                    "shape_ids": [shape_id],
                    "shape_types": [shape_type],
                    "shape_name": shape_name,
                    "shape_index": shape_index,
                    "speaker_notes_included": False,
                }
                if getattr(shape, "has_table", False):
                    headers = [cell.text.strip() for cell in shape.table.rows[0].cells]
                    for row_index, row in enumerate(shape.table.rows, start=1):
                        cells = [cell.text.strip() for cell in row.cells]
                        text = _table_row_text(
                            headers, cells, f"Slide {slide_number} table {shape_index}", row_index
                        )
                        if text:
                            result.append(SourceChunk(text, "powerpoint", {
                                **common, "content_type": "table",
                                "row_start": row_index, "row_end": row_index,
                                "header_context": headers,
                            }))
                elif getattr(shape, "has_text_frame", False):
                    text = "\n".join(
                        paragraph.text.strip() for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    )
                    if text:
                        result.extend(_split(text, "powerpoint", {
                            **common, "content_type": (
                                "chart_label" if shape_type == "chart" else "slide_text"
                            ),
                        }))
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    result.extend(_split(notes, "powerpoint", {
                        "slide_start": slide_number,
                        "slide_end": slide_number,
                        "slide_number": slide_number,
                        "shape_ids": [],
                        "shape_types": ["speaker_notes"],
                        "speaker_notes_included": True,
                        "content_type": "speaker_notes",
                    }))
            except (AttributeError, ValueError):
                pass
        return result
    except Exception as error:
        raise DocumentParseError("The PowerPoint presentation could not be read.") from error


def _xlsx(
    path: Path,
    include_hidden: bool,
    include_very_hidden: bool,
) -> list[SourceChunk]:
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter, range_boundaries

        workbook = load_workbook(path, read_only=False, data_only=False)
        values_workbook = load_workbook(path, read_only=False, data_only=True)
        result: list[SourceChunk] = []
        try:
            if len(workbook.worksheets) > settings.max_workbook_sheets:
                raise DocumentParseError("Workbook exceeds the configured sheet limit.")
            total_rows = sum(sheet.max_row for sheet in workbook.worksheets)
            if total_rows > settings.max_workbook_rows:
                raise DocumentParseError("Workbook exceeds the configured row limit.")
            for worksheet in workbook.worksheets:
                values_sheet = values_workbook[worksheet.title]
                hidden = worksheet.sheet_state != "visible"
                if (
                    worksheet.sheet_state == "hidden" and not include_hidden
                    or worksheet.sheet_state == "veryHidden"
                    and not include_very_hidden
                ):
                    continue
                populated_rows = [
                    row for row in worksheet.iter_rows()
                    if any(cell.value is not None for cell in row)
                ]
                merged_values = _vertical_merged_values(worksheet)
                row_values = [
                    (
                        row[0].row,
                        [
                            merged_values.get((cell.row, cell.column), cell.value)
                            for cell in row
                        ],
                    )
                    for row in populated_rows
                ]
                header_row_number = _header_row_number(row_values)
                header_context = [
                    str(merged_values.get((cell.row, cell.column), cell.value))
                    for cell in (
                        worksheet[header_row_number] if header_row_number else []
                    )
                    if _is_nonempty_cell(merged_values.get((cell.row, cell.column), cell.value))
                ]
                header_by_column = _deduplicated_headers([
                    (
                        cell.column,
                        merged_values.get((cell.row, cell.column), cell.value),
                    )
                    for cell in (worksheet[header_row_number] if header_row_number else [])
                ])
                tables: list[tuple[str, tuple[int, int, int, int]]] = [
                    (table.name, range_boundaries(table.ref))
                    for table in worksheet.tables.values()
                ]
                merged_ranges = [str(value) for value in worksheet.merged_cells.ranges]
                for row in populated_rows:
                    populated = [cell for cell in row if cell.value is not None]
                    if not populated:
                        continue
                    min_column = min(cell.column for cell in populated)
                    max_column = max(cell.column for cell in populated)
                    row_number = populated[0].row
                    cell_range = (
                        f"{get_column_letter(min_column)}{row_number}:"
                        f"{get_column_letter(max_column)}{row_number}"
                    )
                    formulas = {
                        cell.coordinate: str(cell.value)
                        for cell in populated
                        if isinstance(cell.value, str) and cell.value.startswith("=")
                    }
                    values = [
                        (
                            f"{header_by_column.get(cell.column, cell.coordinate)}: "
                            f"{merged_values.get((cell.row, cell.column), cell.value)}"
                            + (
                                f" (cached value: {values_sheet[cell.coordinate].value})"
                                if cell.coordinate in formulas
                                and values_sheet[cell.coordinate].value is not None
                                else ""
                            )
                        )
                        for cell in populated
                    ]
                    table_name = next(
                        (
                            name
                            for name, (table_min_col, table_min_row, table_max_col, table_max_row)
                            in tables
                            if table_min_row <= row_number <= table_max_row
                            and min_column >= table_min_col
                            and max_column <= table_max_col
                        ),
                        None,
                    )
                    result.append(SourceChunk(
                        f"Sheet: {worksheet.title} | " + " | ".join(values),
                        "excel",
                        {
                        "sheet_name": worksheet.title,
                        "hidden_sheet": hidden,
                        "sheet_hidden": hidden,
                        "row_start": row_number,
                        "row_end": row_number,
                        "column_start": get_column_letter(min_column),
                        "column_end": get_column_letter(max_column),
                        "cell_range": cell_range,
                        "table_name": table_name,
                        "header_rows": (
                            [header_row_number] if header_row_number else []
                        ),
                        "header_context": header_context,
                        "merged_ranges": merged_ranges,
                        "formulas": formulas,
                        },
                    ))
        finally:
            workbook.close()
            values_workbook.close()
        return result
    except Exception as error:
        raise DocumentParseError("The Excel workbook could not be read.") from error


def _xls(
    path: Path,
    include_hidden: bool,
    include_very_hidden: bool,
) -> list[SourceChunk]:
    try:
        import xlrd
        from openpyxl.utils import get_column_letter

        workbook = xlrd.open_workbook(path, on_demand=True)
        result: list[SourceChunk] = []
        try:
            for sheet in workbook.sheets():
                hidden = bool(getattr(sheet, "visibility", 0))
                visibility = int(getattr(sheet, "visibility", 0) or 0)
                if (
                    visibility == 1 and not include_hidden
                    or visibility >= 2 and not include_very_hidden
                ):
                    continue
                header_row = next(
                    (
                        row_index + 1
                        for row_index in range(sheet.nrows)
                        if any(sheet.cell_value(row_index, column) not in ("", None)
                               for column in range(sheet.ncols))
                    ),
                    None,
                )
                header_context = (
                    [
                        str(sheet.cell_value(header_row - 1, column))
                        for column in range(sheet.ncols)
                        if sheet.cell_value(header_row - 1, column) not in ("", None)
                    ]
                    if header_row else []
                )
                header_by_column = _deduplicated_headers([
                    (column, sheet.cell_value(header_row - 1, column))
                    for column in range(sheet.ncols)
                ]) if header_row else {}
                for row_index in range(sheet.nrows):
                    populated = [
                        column for column in range(sheet.ncols)
                        if sheet.cell_value(row_index, column) not in ("", None)
                    ]
                    if not populated:
                        continue
                    row_number = row_index + 1
                    column_start = get_column_letter(min(populated) + 1)
                    column_end = get_column_letter(max(populated) + 1)
                    result.append(SourceChunk(
                        f"Sheet: {sheet.name} | " + " | ".join(
                            f"{header_by_column.get(column, get_column_letter(column + 1))}: "
                            f"{sheet.cell_value(row_index, column)}"
                            for column in populated
                        ),
                        "excel",
                        {
                            "sheet_name": sheet.name,
                            "hidden_sheet": hidden,
                            "row_start": row_number,
                            "row_end": row_number,
                            "column_start": column_start,
                            "column_end": column_end,
                            "cell_range": (
                                f"{column_start}{row_number}:"
                                f"{column_end}{row_number}"
                            ),
                            "table_name": None,
                            "header_rows": [header_row] if header_row else [],
                            "header_context": header_context,
                            "merged_ranges": [],
                            "formulas": {},
                        },
                    ))
        finally:
            workbook.release_resources()
        return result
    except Exception as error:
        raise DocumentParseError("The Excel workbook could not be read.") from error


def _docx(path: Path) -> list[SourceChunk]:
    try:
        from docx import Document

        document = Document(str(path))
        result = _docx_paragraph_chunks(document)
        for table_index, table in enumerate(document.tables, start=1):
            headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
            for row_index, row in enumerate(table.rows, start=1):
                cells = [cell.text.strip() for cell in row.cells]
                text = _table_row_text(headers, cells, f"DOCX table {table_index}", row_index)
                if text:
                    result.append(SourceChunk(text, "word", {
                        "table_number": table_index,
                        "row_start": row_index,
                        "row_end": row_index,
                        "header_context": headers,
                        "content_type": "table",
                    }))
        result.extend(_docx_image_chunks(document))
        return result
    except DocumentParseError:
        raise
    except Exception as error:
        raise DocumentParseError("The DOCX file could not be read.") from error


def _csv(path: Path) -> list[SourceChunk]:
    try:
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
            sample = handle.read(8192)
            handle.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample)
            except csv.Error:
                dialect = csv.excel
            rows = [
                (row_number, [str(value).strip() for value in row])
                for row_number, row in enumerate(
                    csv.reader(handle, dialect),
                    start=1,
                )
                if any(str(value).strip() for value in row)
            ]
            headers = rows[0][1] if rows else []
            return [
                SourceChunk(
                    " | ".join(
                        f"{(headers[index] if index < len(headers) and headers[index] else f'Column {index + 1}')}: "
                        f"{value}"
                        for index, value in enumerate(row)
                        if value
                    ),
                    "csv",
                    {"row_start": row_number, "row_end": row_number},
                )
                for row_number, row in rows
            ]
    except Exception as error:
        raise DocumentParseError("The CSV file could not be read.") from error


def _text(path: Path) -> list[SourceChunk]:
    lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    result: list[SourceChunk] = []
    for line_number, value in enumerate(lines, start=1):
        if value.strip():
            result.append(SourceChunk(
                value.strip(), "text",
                {"line_start": line_number, "line_end": line_number},
            ))
    return result


def extract_source_chunks(
    path: Path,
    *,
    include_hidden: bool = True,
    include_very_hidden: bool = False,
) -> list[SourceChunk]:
    started_at = monotonic()
    extension = path.suffix.lower()
    if extension == ".pdf":
        chunks = _pdf(path)
    elif extension == ".pptx":
        chunks = _pptx(path)
    elif extension == ".xlsx":
        chunks = _xlsx(path, include_hidden, include_very_hidden)
    elif extension == ".xls":
        chunks = _xls(path, include_hidden, include_very_hidden)
    elif extension == ".docx":
        chunks = _docx(path)
    elif extension == ".csv":
        chunks = _csv(path)
    elif extension == ".txt":
        chunks = _text(path)
    else:
        text = extract_text(path)
        if extension in IMAGE_EXTENSIONS:
            chunks = [
                SourceChunk(
                    value,
                    "image",
                    {"page_start": 1, "page_end": 1, "part": index, "content_type": "image_ocr"},
                )
                for index, value in enumerate(chunk_image_text(text), start=1)
            ]
        else:
            if extension == ".ppt":
                chunks = _split(text, "powerpoint", {
                    "slide_start": 1,
                    "slide_end": 1,
                    "shape_ids": ["legacy-text-stream"],
                    "shape_types": ["legacy_text"],
                    "speaker_notes_included": False,
                    "content_type": "legacy_text",
                })
            else:
                chunks = _split(text, "text", {
                    "line_start": 1,
                    "line_end": max(1, len(text.splitlines())),
                })
    if not chunks:
        raise DocumentParseError("No readable content was found in the uploaded file.")
    for chunk in chunks:
        validate_source_location(chunk.source_type, chunk.location)
    if monotonic() - started_at > settings.parser_timeout_seconds:
        raise DocumentParseError("Document parsing exceeded the configured time limit.")
    return chunks


def extract_source_metadata(
    path: Path,
    *,
    include_hidden: bool = True,
    include_very_hidden: bool = False,
) -> dict[str, object]:
    """Return workbook-level facts used for administrative status and QA."""
    if path.suffix.lower() not in {".xlsx", ".xls"}:
        return {}
    if path.suffix.lower() == ".xls":
        try:
            import xlrd

            workbook = xlrd.open_workbook(path, on_demand=True)
            try:
                sheets = workbook.sheets()
                supported = [
                    sheet for sheet in sheets
                    if (
                        int(getattr(sheet, "visibility", 0) or 0) == 0
                        or int(getattr(sheet, "visibility", 0) or 0) == 1
                        and include_hidden
                        or int(getattr(sheet, "visibility", 0) or 0) >= 2
                        and include_very_hidden
                    )
                ]
                return {
                    "source_type": "excel",
                    "sheet_count": len(sheets),
                    "sheet_names": [sheet.name for sheet in sheets],
                    "visible_sheet_count": sum(
                        not bool(getattr(sheet, "visibility", 0))
                        for sheet in sheets
                    ),
                    "processed_sheet_names": [sheet.name for sheet in supported],
                    "detected_tables": [],
                    "total_non_empty_rows": sum(
                        1
                        for sheet in supported
                        for row_index in range(sheet.nrows)
                        if any(
                            sheet.cell_value(row_index, column) not in ("", None)
                            for column in range(sheet.ncols)
                        )
                    ),
                }
            finally:
                workbook.release_resources()
        except Exception as error:
            raise DocumentParseError(
                "The Excel workbook metadata could not be read."
            ) from error
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=False, data_only=False)
        try:
            sheets = list(workbook.worksheets)
            supported = [
                sheet for sheet in sheets
                if (
                    sheet.sheet_state == "visible"
                    or sheet.sheet_state == "hidden" and include_hidden
                    or sheet.sheet_state == "veryHidden" and include_very_hidden
                )
            ]
            return {
                "source_type": "excel",
                "sheet_count": len(sheets),
                "sheet_names": [sheet.title for sheet in sheets],
                "visible_sheet_count": sum(
                    sheet.sheet_state == "visible" for sheet in sheets
                ),
                "processed_sheet_names": [sheet.title for sheet in supported],
                "detected_tables": [
                    {"sheet_name": sheet.title, "table_name": table.name}
                    for sheet in supported
                    for table in sheet.tables.values()
                ],
                "total_non_empty_rows": sum(
                    1
                    for sheet in supported
                    for row in sheet.iter_rows()
                    if any(cell.value is not None for cell in row)
                ),
            }
        finally:
            workbook.close()
    except Exception as error:
        raise DocumentParseError(
            "The Excel workbook metadata could not be read."
        ) from error
