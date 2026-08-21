"""Extensible registry of parsers that convert supported uploads to plain text."""

from __future__ import annotations

import csv
import struct
from pathlib import Path
from typing import Protocol

from app.services.image_processor import (
    IMAGE_EXTENSIONS,
    ImageProcessingError,
    extract_image_text,
)


class DocumentParseError(ValueError):
    """Raised when a supported file cannot be converted to readable text."""

    def __init__(self, message: str, code: str = "document_parse_failed") -> None:
        super().__init__(message)
        self.code = code


class DocumentParser(Protocol):
    def extract_text(self, file_path: Path) -> str:
        """Extract plain text from one file."""


class TxtParser:
    def extract_text(self, file_path: Path) -> str:
        return file_path.read_text(encoding="utf-8", errors="replace")


class PdfParser:
    def extract_text(self, file_path: Path) -> str:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as error:
            raise DocumentParseError("The PDF file could not be read.") from error


class DocxParser:
    def extract_text(self, file_path: Path) -> str:
        try:
            from docx import Document

            document = Document(str(file_path))
            return "\n".join(paragraph.text for paragraph in document.paragraphs)
        except Exception as error:
            raise DocumentParseError("The DOCX file could not be read.") from error


def _cell_text(value: object) -> str:
    if value is None:
        return ""
    return str(value).replace("\r\n", "\n").replace("\r", "\n").strip()


class ExcelParser:
    def extract_text(self, file_path: Path) -> str:
        if file_path.suffix.lower() == ".xlsx":
            return self._extract_xlsx(file_path)
        return self._extract_xls(file_path)

    def _extract_xlsx(self, file_path: Path) -> str:
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(file_path, read_only=True, data_only=True)
            try:
                output = [f"Workbook: {file_path.stem}"]
                for worksheet in workbook.worksheets:
                    output.append(f"Sheet: {worksheet.title}")
                    for row in worksheet.iter_rows(values_only=True):
                        values = [_cell_text(value) for value in row]
                        if any(values):
                            output.append("\t".join(values).rstrip())
                return "\n".join(output)
            finally:
                workbook.close()
        except DocumentParseError:
            raise
        except Exception as error:
            raise DocumentParseError("The Excel workbook could not be read.") from error

    def _extract_xls(self, file_path: Path) -> str:
        try:
            import xlrd

            workbook = xlrd.open_workbook(str(file_path), on_demand=True)
            try:
                output = [f"Workbook: {file_path.stem}"]
                for worksheet in workbook.sheets():
                    output.append(f"Sheet: {worksheet.name}")
                    for row_index in range(worksheet.nrows):
                        values = [_cell_text(worksheet.cell_value(row_index, column)) for column in range(worksheet.ncols)]
                        if any(values):
                            output.append("\t".join(values).rstrip())
                return "\n".join(output)
            finally:
                workbook.release_resources()
        except Exception as error:
            raise DocumentParseError("The legacy Excel workbook could not be read.") from error


class CsvParser:
    def extract_text(self, file_path: Path) -> str:
        try:
            output = [f"CSV: {file_path.stem}"]
            with file_path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
                sample = handle.read(8192)
                handle.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample)
                except csv.Error:
                    dialect = csv.excel
                for row in csv.reader(handle, dialect):
                    values = [_cell_text(value) for value in row]
                    if any(values):
                        output.append("\t".join(values).rstrip())
            return "\n".join(output)
        except Exception as error:
            raise DocumentParseError("The CSV file could not be read.") from error


def _shape_text(shape) -> list[str]:
    if getattr(shape, "has_table", False):
        return ["\t".join(_cell_text(cell.text) for cell in row.cells).rstrip() for row in shape.table.rows]
    if getattr(shape, "has_text_frame", False):
        return [paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
    return []


def _extract_legacy_ppt_text(data: bytes) -> list[str]:
    """Extract PowerPoint text atoms from the legacy OLE record stream."""
    text_record_types = {4000: "utf-16-le", 4008: "cp1252", 4026: "utf-16-le"}
    output: list[str] = []

    def walk(start: int, end: int) -> None:
        offset = start
        while offset + 8 <= end:
            version_instance, record_type, length = struct.unpack_from("<HHI", data, offset)
            payload_start = offset + 8
            payload_end = payload_start + length
            if payload_end > end or payload_end < payload_start:
                break
            record_version = version_instance & 0xF
            if record_type in text_record_types:
                text = data[payload_start:payload_end].decode(text_record_types[record_type], errors="ignore")
                text = text.replace("\x00", "").strip()
                if text and (not output or output[-1] != text):
                    output.append(text)
            elif record_version == 0xF:
                walk(payload_start, payload_end)
            offset = payload_end

    walk(0, len(data))
    return output


class PowerPointParser:
    def extract_text(self, file_path: Path) -> str:
        if file_path.suffix.lower() == ".pptx":
            return self._extract_pptx(file_path)
        return self._extract_ppt(file_path)

    def _extract_pptx(self, file_path: Path) -> str:
        try:
            from pptx import Presentation

            presentation = Presentation(str(file_path))
            output = [f"Presentation: {file_path.stem}"]
            for index, slide in enumerate(presentation.slides, start=1):
                output.append(f"Slide {index}")
                for shape in slide.shapes:
                    output.extend(text for text in _shape_text(shape) if text)
                try:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        output.extend(("Speaker notes", notes))
                except (AttributeError, ValueError):
                    pass
            return "\n".join(output)
        except DocumentParseError:
            raise
        except Exception as error:
            raise DocumentParseError("The PowerPoint presentation could not be read.") from error

    def _extract_ppt(self, file_path: Path) -> str:
        try:
            import olefile

            with olefile.OleFileIO(str(file_path)) as presentation:
                stream_name = "PowerPoint Document"
                if not presentation.exists(stream_name):
                    raise DocumentParseError("The legacy PowerPoint file has no readable presentation stream.")
                values = _extract_legacy_ppt_text(presentation.openstream(stream_name).read())
            return "\n".join([f"Presentation: {file_path.stem}", *values])
        except DocumentParseError:
            raise
        except Exception as error:
            raise DocumentParseError("The legacy PowerPoint presentation could not be read.") from error


class OcrParser:
    def extract_text(self, file_path: Path) -> str:
        try:
            return extract_image_text(file_path)
        except ImageProcessingError as error:
            raise DocumentParseError(str(error), code=error.code) from error


PARSER_REGISTRY: dict[str, DocumentParser] = {}


def register_parser(extensions: tuple[str, ...], parser: DocumentParser) -> None:
    for extension in extensions:
        PARSER_REGISTRY[extension.lower()] = parser


register_parser((".txt",), TxtParser())
register_parser((".pdf",), PdfParser())
register_parser((".docx",), DocxParser())
register_parser((".xlsx", ".xls"), ExcelParser())
register_parser((".csv",), CsvParser())
register_parser((".pptx", ".ppt"), PowerPointParser())
register_parser(tuple(IMAGE_EXTENSIONS), OcrParser())

SUPPORTED_EXTENSIONS = frozenset(PARSER_REGISTRY)


def extract_text(file_path: Path) -> str:
    """Dispatch extraction through the registered parser for the file extension."""
    parser = PARSER_REGISTRY.get(file_path.suffix.lower())
    if parser is None:
        raise DocumentParseError("Unsupported document type.")
    return parser.extract_text(file_path)
