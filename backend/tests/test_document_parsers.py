"""Parser-registry tests for tabular, presentation, and OCR formats."""

from __future__ import annotations

import struct
import tempfile
import unittest
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from app.services.document_loader import (
    DocumentParseError,
    ExcelParser,
    OcrParser,
    PARSER_REGISTRY,
    PowerPointParser,
    SUPPORTED_EXTENSIONS,
    _extract_legacy_ppt_text,
    extract_text,
)
from app.services.source_extraction import (
    extract_source_chunks,
    extract_source_metadata,
    validate_source_location,
)
from app.services.image_processor.image_parser import ImageProcessingError
from app.services.image_processor.ocr import (
    OCR_FAILED_MESSAGE,
    OCR_TIMEOUT_MESSAGE,
    OCR_UNAVAILABLE_MESSAGE,
    OcrResult,
    ocr_health,
    require_ocr_ready_for_startup,
)
from app.services.pdf_layout import PdfPageText


class DocumentParserTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def test_registry_contains_every_required_extension(self):
        required = {".txt", ".pdf", ".docx", ".xlsx", ".xls", ".csv", ".pptx", ".ppt", ".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".webp"}
        self.assertEqual(SUPPORTED_EXTENSIONS, required)
        self.assertIsInstance(PARSER_REGISTRY[".xlsx"], ExcelParser)
        self.assertIsInstance(PARSER_REGISTRY[".ppt"], PowerPointParser)
        self.assertIsInstance(PARSER_REGISTRY[".webp"], OcrParser)

    def test_csv_parser_extracts_rows(self):
        path = self.root / "employees.csv"
        path.write_text("Employee ID,Name,Department\n1001,John,HR\n1002,Alice,Finance", encoding="utf-8")
        text = extract_text(path)
        self.assertIn("CSV: employees", text)
        self.assertIn("Employee ID\tName\tDepartment", text)
        self.assertIn("1002\tAlice\tFinance", text)
        chunks = extract_source_chunks(path)
        self.assertIn("Employee ID: 1002", chunks[2].text)
        self.assertIn("Name: Alice", chunks[2].text)
        self.assertEqual(chunks[2].location["row_start"], 3)

    def test_xlsx_parser_extracts_workbook_sheets_and_cells(self):
        from openpyxl import Workbook

        path = self.root / "employees.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Employees"
        sheet.append(["Employee ID", "Name", "Department"])
        sheet.append([1001, "John", "HR"])
        workbook.create_sheet("Summary").append(["Total", 1])
        workbook.save(path)
        workbook.close()
        text = extract_text(path)
        self.assertIn("Workbook: employees", text)
        self.assertIn("Sheet: Employees", text)
        self.assertIn("1001\tJohn\tHR", text)
        self.assertIn("Sheet: Summary", text)

    def test_pptx_parser_extracts_slide_text_and_tables(self):
        from pptx import Presentation
        from pptx.util import Inches

        path = self.root / "quarterly.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Quarterly Sales"
        slide.placeholders[1].text = "Revenue increased by 15%"
        table = slide.shapes.add_table(1, 2, Inches(1), Inches(4), Inches(6), Inches(1)).table
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Revenue"
        presentation.save(path)
        text = extract_text(path)
        self.assertIn("Presentation: quarterly", text)
        self.assertIn("Slide 1", text)
        self.assertIn("Quarterly Sales", text)
        self.assertIn("Revenue increased by 15%", text)
        self.assertIn("Region\tRevenue", text)

    def test_pdf_chunks_retain_exact_page_numbers(self):
        path = self.root / "policy.pdf"
        path.write_bytes(b"%PDF-test")
        pages = [
            SimpleNamespace(extract_text=lambda: "First page"),
            SimpleNamespace(extract_text=lambda: "Second page"),
        ]
        with patch("pypdf.PdfReader", return_value=SimpleNamespace(pages=pages)):
            chunks = extract_source_chunks(path)
        self.assertEqual(
            [chunk.location["page_start"] for chunk in chunks],
            [1, 2],
        )
        self.assertTrue(all(chunk.source_type == "pdf" for chunk in chunks))

    def test_pptx_chunks_retain_slide_shape_and_table_row(self):
        from pptx import Presentation
        from pptx.util import Inches

        path = self.root / "located.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Located title"
        table = slide.shapes.add_table(
            2, 2, Inches(1), Inches(4), Inches(6), Inches(1)
        ).table
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Revenue"
        table.cell(1, 0).text = "West"
        table.cell(1, 1).text = "120"
        presentation.save(path)
        chunks = extract_source_chunks(path)
        title = next(chunk for chunk in chunks if "Located title" in chunk.text)
        table_chunk = next(chunk for chunk in chunks if "West" in chunk.text)
        self.assertEqual(title.location["slide_number"], 1)
        self.assertEqual(title.location["slide_start"], 1)
        self.assertEqual(title.location["slide_end"], 1)
        self.assertIn("shape_index", title.location)
        self.assertEqual(title.location["shape_ids"], ["shape-1"])
        self.assertFalse(title.location["speaker_notes_included"])
        self.assertIn("Region: West", table_chunk.text)
        self.assertIn("Revenue: 120", table_chunk.text)
        self.assertEqual(table_chunk.location["row_start"], 2)
        self.assertEqual(table_chunk.location["header_context"], ["Region", "Revenue"])
        self.assertEqual(table_chunk.location["content_type"], "table")

    def test_docx_table_rows_include_headers_and_table_provenance(self):
        """DOCX data rows remain structured when header and values are separate cells."""
        from docx import Document

        path = self.root / "regional.docx"
        document = Document()
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Revenue"
        table.cell(1, 0).text = "West"
        table.cell(1, 1).text = "120"
        document.save(path)

        chunks = extract_source_chunks(path)
        data = next(chunk for chunk in chunks if "Region: West" in chunk.text)
        self.assertIn("Revenue: 120", data.text)
        self.assertEqual(data.location["table_number"], 1)
        self.assertEqual(data.location["row_start"], 2)
        self.assertEqual(data.location["header_context"], ["Region", "Revenue"])

    def test_docx_short_paragraphs_are_batched_into_retrieval_chunks(self):
        """Many short DOCX paragraphs should not create one retrieval chunk each."""
        from docx import Document

        path = self.root / "large-prose.docx"
        document = Document()
        for index in range(600):
            document.add_paragraph(f"Requirement {index}: verify account handling.")
        document.save(path)

        chunks = extract_source_chunks(path)

        self.assertLess(len(chunks), 20)
        self.assertEqual(chunks[0].source_type, "word")
        self.assertEqual(chunks[0].location["paragraph_start"], 1)
        self.assertGreater(chunks[0].location["paragraph_end"], 1)
        self.assertIn("Requirement 0", chunks[0].text)
        self.assertIn("Requirement 599", chunks[-1].text)

    def test_docx_headings_split_batched_prose_sections(self):
        """Section-heading questions should retrieve the matching DOCX section."""
        from docx import Document

        path = self.root / "sectioned.docx"
        document = Document()
        document.add_heading("Purpose", level=2)
        document.add_paragraph("The document purpose is validation support.")
        document.add_heading("Product Scope", level=2)
        document.add_paragraph("The product scope includes bonus reports and exports.")
        document.save(path)

        chunks = extract_source_chunks(path)
        purpose = next(chunk for chunk in chunks if "Purpose" in chunk.text)
        scope = next(chunk for chunk in chunks if "Product Scope" in chunk.text)

        self.assertNotEqual(purpose.location["paragraph_start"], scope.location["paragraph_start"])
        self.assertIn("validation support", purpose.text)
        self.assertNotIn("Product Scope", purpose.text)
        self.assertIn("bonus reports and exports", scope.text)

    def test_docx_embedded_image_ocr_creates_image_chunk_with_provenance(self):
        """DOCX embedded media should contribute OCR chunks without losing text/table chunks."""
        from docx import Document
        from PIL import Image

        image_path = self.root / "embedded.png"
        Image.new("RGB", (120, 40), "white").save(image_path)
        path = self.root / "embedded.docx"
        document = Document()
        document.add_heading("Project Alpha", level=1)
        document.add_paragraph("Project Alpha paragraph evidence.")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Project"
        table.cell(0, 1).text = "Current rate"
        table.cell(1, 0).text = "Alpha"
        table.cell(1, 1).text = "17.5"
        document.add_picture(str(image_path))
        document.save(path)

        with patch("app.services.source_extraction.extract_image_text", return_value="OCR text:\nEmbedded approval code is IMG-42."):
            chunks = extract_source_chunks(path)

        self.assertTrue(any("Project Alpha paragraph evidence" in chunk.text for chunk in chunks))
        self.assertTrue(any("Current rate: 17.5" in chunk.text for chunk in chunks))
        image = next(chunk for chunk in chunks if "IMG-42" in chunk.text)
        self.assertEqual(image.source_type, "word")
        self.assertEqual(image.location["content_type"], "image_ocr")
        self.assertIn("relationship_id", image.location)
        self.assertIn("image_index", image.location)

    def test_image_only_docx_ocr_creates_image_chunk(self):
        """Image-only DOCX files should index embedded OCR instead of being rejected."""
        from docx import Document
        from PIL import Image

        image_path = self.root / "only-image.png"
        Image.new("RGB", (120, 40), "white").save(image_path)
        path = self.root / "image-only.docx"
        document = Document()
        document.add_picture(str(image_path))
        document.save(path)

        with patch("app.services.source_extraction.extract_image_text", return_value="OCR text:\nImage-only approval IMG-77."):
            chunks = extract_source_chunks(path)

        self.assertEqual(len(chunks), 1)
        self.assertIn("IMG-77", chunks[0].text)
        self.assertEqual(chunks[0].location["content_type"], "image_ocr")

    def test_docx_embedded_ocr_dependency_failure_is_preserved(self):
        """OCR infrastructure failures must not become generic no-readable-content errors."""
        from docx import Document
        from PIL import Image

        image_path = self.root / "dependency.png"
        Image.new("RGB", (120, 40), "white").save(image_path)
        path = self.root / "dependency.docx"
        document = Document()
        document.add_picture(str(image_path))
        document.save(path)

        with patch(
            "app.services.source_extraction.extract_image_text",
            side_effect=ImageProcessingError(OCR_UNAVAILABLE_MESSAGE, code="ocr_unavailable"),
        ):
            with self.assertRaises(DocumentParseError) as raised:
                extract_source_chunks(path)

        self.assertEqual(raised.exception.code, "ocr_unavailable")
        self.assertEqual(str(raised.exception), OCR_UNAVAILABLE_MESSAGE)
        self.assertNotIn("No readable content", str(raised.exception))

    def test_docx_mixed_text_survives_corrupt_embedded_image(self):
        """Malformed embedded images are skipped when native DOCX text is readable."""
        from app.services.source_extraction import _image_chunks_from_bytes

        chunks = _image_chunks_from_bytes(
            blob=b"not an image",
            filename_hint="broken.png",
            source_type="word",
            location={"content_type": "image_ocr"},
        )
        self.assertEqual(chunks, [])

    def test_pdf_scanned_page_ocr_creates_page_image_chunk(self):
        """A PDF page with only an image should still yield OCR text with page provenance."""
        from PIL import Image

        buffer = BytesIO()
        Image.new("RGB", (120, 40), "white").save(buffer, format="PNG")
        path = self.root / "scanned.pdf"
        path.write_bytes(b"%PDF-test")
        fake_image = SimpleNamespace(data=buffer.getvalue(), name="scan.png")
        fake_page = SimpleNamespace(images=[fake_image])
        with (
            patch("app.services.source_extraction.extract_pdf_page_texts", return_value=[]),
            patch("pypdf.PdfReader", return_value=SimpleNamespace(pages=[fake_page])),
            patch("app.services.source_extraction.extract_image_text", return_value="OCR text:\nScanned approval code is PDF-42."),
        ):
            chunks = extract_source_chunks(path)

        scanned = next(chunk for chunk in chunks if "PDF-42" in chunk.text)
        self.assertEqual(scanned.source_type, "pdf")
        self.assertEqual(scanned.location["page_start"], 1)
        self.assertEqual(scanned.location["content_type"], "image_ocr")

    def test_pdf_native_and_image_ocr_are_both_preserved_when_distinct(self):
        """Native text pages and distinct embedded-image OCR stay complementary."""
        from PIL import Image

        buffer = BytesIO()
        Image.new("RGB", (120, 40), "white").save(buffer, format="PNG")
        path = self.root / "native-image.pdf"
        path.write_bytes(b"%PDF-test")
        fake_page = SimpleNamespace(images=[SimpleNamespace(data=buffer.getvalue(), name="embedded.png")])
        with (
            patch("app.services.source_extraction.extract_pdf_page_texts", return_value=[PdfPageText(1, "Native text says base rate is 12.")]),
            patch("pypdf.PdfReader", return_value=SimpleNamespace(pages=[fake_page])),
            patch("app.services.source_extraction.extract_image_text", return_value="OCR text:\nImage note says approval is required."),
        ):
            chunks = extract_source_chunks(path)

        self.assertTrue(any("base rate is 12" in chunk.text for chunk in chunks))
        image = next(chunk for chunk in chunks if "approval is required" in chunk.text)
        self.assertEqual(image.location["page_start"], 1)
        self.assertEqual(image.location["content_type"], "image_ocr")

    def test_pdf_ocr_duplicate_text_is_suppressed(self):
        """Native/OCR duplicate suppression avoids indexing the same text twice."""
        from PIL import Image

        buffer = BytesIO()
        Image.new("RGB", (120, 40), "white").save(buffer, format="PNG")
        path = self.root / "duplicate.pdf"
        path.write_bytes(b"%PDF-test")
        native = "Alpha rate requires approval from finance team before release."
        fake_page = SimpleNamespace(images=[SimpleNamespace(data=buffer.getvalue(), name="duplicate.png")])
        with (
            patch("app.services.source_extraction.extract_pdf_page_texts", return_value=[PdfPageText(1, native)]),
            patch("pypdf.PdfReader", return_value=SimpleNamespace(pages=[fake_page])),
            patch("app.services.source_extraction.extract_image_text", return_value=f"OCR text:\n{native}"),
        ):
            chunks = extract_source_chunks(path)

        self.assertEqual(sum("Alpha rate requires approval" in chunk.text for chunk in chunks), 1)

    def test_embedded_ocr_limit_can_reject_otherwise_unreadable_scanned_pdf(self):
        """A disabled embedded-OCR limit keeps unreadable scanned PDFs safely rejected."""
        from PIL import Image

        buffer = BytesIO()
        Image.new("RGB", (120, 40), "white").save(buffer, format="PNG")
        path = self.root / "limited.pdf"
        path.write_bytes(b"%PDF-test")
        fake_page = SimpleNamespace(images=[SimpleNamespace(data=buffer.getvalue(), name="scan.png")])
        with (
            patch("app.services.source_extraction.extract_pdf_page_texts", return_value=[]),
            patch("pypdf.PdfReader", return_value=SimpleNamespace(pages=[fake_page])),
            patch("app.services.source_extraction.extract_image_text", return_value="OCR text should be blocked."),
            patch("app.services.source_extraction.settings.embedded_ocr_max_images_per_document", 0),
        ):
            with self.assertRaisesRegex(DocumentParseError, "No readable content"):
                extract_source_chunks(path)

    def test_pdf_table_rows_include_headers_and_page_provenance(self):
        """Layout-preserved PDF tables do not leave values in headerless prose chunks."""
        path = self.root / "regional.pdf"
        path.write_bytes(b"%PDF-test")
        with patch(
            "app.services.source_extraction.extract_pdf_page_texts",
            return_value=[PdfPageText(3, "Region | Revenue\nWest | 120\nEast | 140")],
        ):
            chunks = extract_source_chunks(path)

        west = next(chunk for chunk in chunks if "Region: West" in chunk.text)
        self.assertIn("Revenue: 120", west.text)
        self.assertEqual(west.location["page_start"], 3)
        self.assertEqual(west.location["table_number"], 1)
        self.assertEqual(west.location["header_context"], ["Region", "Revenue"])

    def test_xlsx_chunks_retain_sheet_cell_range_table_and_formula(self):
        from openpyxl import Workbook
        from openpyxl.worksheet.table import Table

        path = self.root / "located.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Revenue"
        sheet.append(["Amount", "Tax"])
        sheet.append([100, "=A2*0.1"])
        sheet.add_table(Table(displayName="RevenueTable", ref="A1:B2"))
        workbook.save(path)
        workbook.close()
        chunks = extract_source_chunks(path)
        formula_row = next(chunk for chunk in chunks if "=A2*0.1" in chunk.text)
        self.assertIn("Sheet: Revenue", formula_row.text)
        self.assertIn("Amount: 100", formula_row.text)
        self.assertIn("Tax: =A2*0.1", formula_row.text)
        self.assertEqual(formula_row.location["sheet_name"], "Revenue")
        self.assertEqual(formula_row.location["cell_range"], "A2:B2")
        self.assertEqual(formula_row.location["table_name"], "RevenueTable")
        self.assertEqual(formula_row.location["formulas"], {"B2": "=A2*0.1"})
        self.assertEqual(formula_row.location["header_rows"], [1])
        self.assertEqual(formula_row.location["header_context"], ["Amount", "Tax"])
        self.assertFalse(formula_row.location["hidden_sheet"])
        metadata = extract_source_metadata(path)
        self.assertEqual(metadata["sheet_count"], 1)
        self.assertEqual(metadata["visible_sheet_count"], 1)
        self.assertEqual(metadata["total_non_empty_rows"], 2)
        self.assertEqual(
            metadata["detected_tables"],
            [{"sheet_name": "Revenue", "table_name": "RevenueTable"}],
        )

    def test_source_location_validation_rejects_incomplete_citations(self):
        with self.assertRaisesRegex(DocumentParseError, "missing page_end"):
            validate_source_location("pdf", {"page_start": 1})

    def test_legacy_ppt_text_atom_extraction(self):
        payload = "Legacy slide text".encode("utf-16-le")
        record = struct.pack("<HHI", 0, 4000, len(payload)) + payload
        self.assertEqual(_extract_legacy_ppt_text(record), ["Legacy slide text"])

    def test_ocr_parser_uses_all_image_frames(self):
        from PIL import Image

        path = self.root / "scan.png"
        Image.new("RGB", (30, 20), "white").save(path)
        with (
            patch("pytesseract.image_to_string", return_value="Invoice total 42") as ocr,
            patch(
                "app.services.image_processor.image_parser.vision_is_configured",
                return_value=True,
            ),
            patch(
                "app.services.image_processor.image_parser.describe_image",
                return_value="A scanned invoice with a visible total.",
            ),
        ):
            text = extract_text(path)
        self.assertIn("Image: scan", text)
        self.assertIn("Invoice total 42", text)
        self.assertIn("A scanned invoice with a visible total.", text)
        ocr.assert_called_once()

    def test_missing_tesseract_returns_clear_error(self):
        import pytesseract
        from PIL import Image

        path = self.root / "scan.jpg"
        Image.new("RGB", (30, 20), "white").save(path)
        with (
            patch("pytesseract.image_to_string", side_effect=pytesseract.TesseractNotFoundError()),
            patch(
                "app.services.image_processor.image_parser.vision_is_configured",
                return_value=False,
            ),
        ):
            with self.assertRaisesRegex(DocumentParseError, "OCR is currently unavailable"):
                extract_text(path)

    def test_direct_png_and_jpeg_ocr_chunks_include_image_metadata(self):
        """Direct image ingestion should retain OCR source metadata for PNG and JPEG."""
        from PIL import Image

        for filename, image_format in (("scan.png", "PNG"), ("scan.jpg", "JPEG")):
            path = self.root / filename
            Image.new("RGB", (30, 20), "white").save(path)
            with (
                patch(
                    "app.services.image_processor.image_parser.extract_ocr",
                    return_value=OcrResult("Readable approval text", image_format, 30, 20, 1),
                ),
                patch(
                    "app.services.image_processor.image_parser.vision_is_configured",
                    return_value=False,
                ),
            ):
                chunks = extract_source_chunks(path)
            self.assertTrue(chunks)
            self.assertEqual(chunks[0].source_type, "image")
            self.assertEqual(chunks[0].location["content_type"], "image_ocr")

    def test_invalid_tesseract_cmd_returns_sanitized_unavailable_error(self):
        import pytesseract
        from PIL import Image

        path = self.root / "invalid-cmd.png"
        Image.new("RGB", (30, 20), "white").save(path)
        with (
            patch("app.services.image_processor.ocr.settings.tesseract_cmd", "C:/missing/tesseract.exe"),
            patch("pytesseract.image_to_string", side_effect=pytesseract.TesseractNotFoundError()),
            patch("app.services.image_processor.image_parser.vision_is_configured", return_value=False),
        ):
            with self.assertRaises(DocumentParseError) as raised:
                extract_text(path)
        self.assertEqual(raised.exception.code, "ocr_unavailable")
        self.assertEqual(str(raised.exception), OCR_UNAVAILABLE_MESSAGE)
        self.assertNotIn("C:/missing", str(raised.exception))

    def test_ocr_timeout_and_no_text_have_specific_messages(self):
        from PIL import Image

        path = self.root / "timeout.png"
        Image.new("RGB", (30, 20), "white").save(path)
        with (
            patch("pytesseract.image_to_string", side_effect=RuntimeError("timeout")),
            patch("app.services.image_processor.image_parser.vision_is_configured", return_value=False),
        ):
            with self.assertRaises(DocumentParseError) as timed_out:
                extract_text(path)
        self.assertEqual(timed_out.exception.code, "ocr_timeout")
        self.assertEqual(str(timed_out.exception), OCR_TIMEOUT_MESSAGE)

        with (
            patch(
                "app.services.image_processor.image_parser.extract_ocr",
                return_value=OcrResult("", "PNG", 30, 20, 1),
            ),
            patch("app.services.image_processor.image_parser.vision_is_configured", return_value=False),
        ):
            with self.assertRaises(DocumentParseError) as no_text:
                extract_text(path)
        self.assertEqual(no_text.exception.code, "ocr_no_text")
        self.assertEqual(str(no_text.exception), "No readable text was detected in the image.")

    def test_ocr_health_and_production_startup_are_safe(self):
        """Health exposes only ready/unavailable and production startup fails fast."""
        with (
            patch("pytesseract.get_tesseract_version", return_value="5.3.0"),
            patch("pytesseract.get_languages", return_value=["eng"]),
            patch("app.services.image_processor.ocr.shutil.which", return_value="tesseract"),
        ):
            self.assertEqual(ocr_health(), {"status": "ready"})

        with (
            patch("pytesseract.get_tesseract_version", side_effect=RuntimeError("path C:/secret failed")),
            patch("app.services.image_processor.ocr.settings.app_environment", "production"),
        ):
            self.assertEqual(ocr_health(), {"status": "unavailable"})
            with self.assertRaisesRegex(RuntimeError, "OCR dependency is unavailable"):
                require_ocr_ready_for_startup()

    def test_corrupt_excel_and_powerpoint_return_format_errors(self):
        xlsx = self.root / "broken.xlsx"
        pptx = self.root / "broken.pptx"
        xlsx.write_bytes(b"not a workbook")
        pptx.write_bytes(b"not a presentation")
        with self.assertRaisesRegex(DocumentParseError, "Excel workbook"):
            extract_text(xlsx)
        with self.assertRaisesRegex(DocumentParseError, "PowerPoint presentation"):
            extract_text(pptx)

    def test_unsupported_extension_is_rejected(self):
        path = self.root / "document.exe"
        path.write_bytes(b"unsafe")
        with self.assertRaisesRegex(DocumentParseError, "Unsupported document type"):
            extract_text(path)


if __name__ == "__main__":
    unittest.main()
